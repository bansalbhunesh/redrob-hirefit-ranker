#!/usr/bin/env python3
"""One-off: refresh deck numbers after the 2026-06-11 consensus calibration roll.

Post-roll judge-1 metrics on the calibrated submission: NDCG@50 0.8734 -> 0.9245,
composite 0.8959 -> 0.9112; independent composite 0.8811 -> 0.8863. The
LambdaMART panel is updated to the four-negatives ledger + the adopted
calibration; the measured-runtime row picks up the fresh-build 124.7 s serial
datapoint. NDCG@10 / MAP / P@10 cards are unchanged by the roll (top-10
untouched).

Usage: python scripts/roll_deck_numbers_2.py docs/HireFit_Ranker_Redrob_POLISHED.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

RUN_PATCHES = [
    ("0.873", "0.924"),        # slide 8 card: judge-1 NDCG@50
    ("(0.896)", "(0.911)"),    # slide 8 panel: judge composite
    ("(0.881)", "(0.886)"),    # slide 8 panel: independent composite
    ("80–122 s", "80–125 s"),
    ("80-122 s", "80-125 s"),
]

DISCIPLINE_CARD = ("0.896", "0.911")  # SHIPPED (LLM-JUDGE) stat card

NEW_RIVAL_BODY = (
    "A LambdaMART challenger built on our own 28 features plus recovered "
    "generator structure lost its pre-registered gate: 0.900 vs 0.906. One of "
    "four measured negatives — embeddings, learned-LR, the challenger, and a "
    "declined availability hedge. The single adopted change: an 8-swap "
    "consensus calibration pass that survived held-out validation "
    "(+0.009–0.011)."
)


def patch_runs(prs) -> int:
    hits = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    for old, new in RUN_PATCHES:
                        if old in t:
                            t = t.replace(old, new)
                    if t != run.text:
                        run.text = t
                        hits += 1
    return hits


def main(path: str) -> None:
    prs = Presentation(path)

    # Discipline slide (index 3 since the polish pass): SHIPPED card + rival panel.
    card = panel = 0
    for shape in prs.slides[3].shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == DISCIPLINE_CARD[0]:
            runs = shape.text_frame.paragraphs[0].runs
            if runs and runs[0].font.size and runs[0].font.size.pt >= 20:
                runs[0].text = DISCIPLINE_CARD[1]
                card += 1
        elif text.startswith("A LambdaMART ranker"):
            runs = shape.text_frame.paragraphs[0].runs
            runs[0].text = NEW_RIVAL_BODY
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
            panel += 1
    print(f"discipline card patched: {card}; rival panel rewritten: {panel}")

    print(f"run-level patches: {patch_runs(prs)}")
    prs.save(path)
    print(f"saved {Path(path).name}")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "docs/HireFit_Ranker_Redrob_POLISHED.pptx")
