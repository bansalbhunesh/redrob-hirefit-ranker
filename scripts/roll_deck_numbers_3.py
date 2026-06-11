#!/usr/bin/env python3
"""One-off: audit-v2 deck repairs (2026-06-11, docs/deck only — submission frozen).

Every patched number traces to a committed artifact:

- Slide 3 "penalize overshoot" claim: contradicted by the shipped code
  (src/redrob_ranker/features.py, yoe_fit_score — yoe >= band center scores
  1.0). Reworded to what the code does: juniors demoted hard (Gaussian below
  band + 0.55 cut under 3y); senior depth read from trajectory/title, not
  years arithmetic.
- Slide 4 "0.926 FLOOR-OFF — DECLINED": the number is the floor-1.00 LLM-side
  composite at 67% coverage that docs/sensitivity_sweep.md's own comparability
  guard excludes as selection-biased — the card label now says so.
- Slide 4 honest-framing note "demote all 12": miscount. Per
  docs/generator_forensics.md the gold strata split 18 in / 11 out, and the
  gold-clothed honeypot (CAND_0093547) is one OF the 11. Card 11/11 stands;
  the note is corrected to 11.
- Slide 8 "570+": untraceable to any committed artifact (repo-wide search,
  audit v2). Replaced with a verified stat: the pool's 12 non-target strata
  hold >= 4,800 wrong-role profiles (12 x 405-504, generator_forensics.md)
  and a summary-template scan of the shipped top-100 places all 100 ranks in
  the four target strata — zero wrong-role entrants.
- Slide 9 honest-framing strip: adds the post-adoption judge-3 re-test of the
  calibration (+0.0124, 6 confirm / 1 tie / 1 contested —
  docs/llm_judge_eval_3.md), replacing reliance on the low-power crossover.
- Slide 12 "<50 ms": hardcoded with no artifact; replaced with the measured
  ~3 ms (docs/finale_runbook.md, 5-request sample 2.7-3.8 ms).

Usage: python scripts/roll_deck_numbers_3.py [docs/HireFit_Ranker_Redrob_POLISHED.pptx]
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

# (slide_index, exact stripped shape text to find, replacement paragraph text)
SHAPE_REWRITES = [
    (
        2,
        "Right seniority  —  5–9 years of judgment — penalize overshoot and juniors.",
        "Right seniority  —  anchored on the JD's 5–9 band — juniors demoted "
        "hard; senior depth read from trajectory, not years arithmetic.",
    ),
    (
        3,
        "FLOOR-OFF — DECLINED",
        "FLOOR-OFF — BIASED SAMPLE, DECLINED",
    ),
    (
        3,
        "Honest framing:  the pool plants 11 perfect-on-paper-but-unavailable seniors "
        "(response rates to 0.07) and one honeypot in gold clothing — the guardrails "
        "demote all 12; an LLM judge missed one of them.",
        "Honest framing:  the pool plants 11 perfect-on-paper-but-unavailable gold "
        "profiles — one a honeypot in gold clothing — the guardrails demote all "
        "11; an LLM judge missed one of them.",
    ),
    (
        7,
        "570+",
        "4,800+",
    ),
    (
        7,
        "Lowest rank a non-target-title profile with a high AI-skill match reaches. "
        "None enter the top 100.",
        "Wrong-role profiles across the pool's 12 non-target strata — zero reach "
        "the top-100.",
    ),
    (
        8,
        "Honest framing:  this is a development-time proxy on a 249-candidate sample, "
        "not the hidden challenge score. It tells us the top of the list is real — "
        "it is not a leaderboard claim.",
        "Honest framing:  a development-time proxy on a 249-candidate sample, not the "
        "hidden challenge score. The one post-freeze change (8-swap calibration) was "
        "re-tested by a judge family added after adoption: +0.0124, 6 confirm / 1 tie "
        "/ 1 contested (docs/llm_judge_eval_3.md).",
    ),
    (
        11,
        "<50 ms",
        "~3 ms",
    ),
    (
        11,
        "PER-CANDIDATE AUDIT API",
        "AUDIT API, MEASURED",
    ),
]


def rewrite_shape(slide, old: str, new: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() != old:
            continue
        # Collapse the matched text into the first run of the first non-empty
        # paragraph (the deck builder styles whole stat cards / strips per run,
        # so keeping run 0's formatting preserves the look).
        for para in shape.text_frame.paragraphs:
            if not para.runs:
                continue
            para.runs[0].text = new
            for r in para.runs[1:]:
                r._r.getparent().remove(r._r)
            # Wipe any remaining paragraphs' runs (multi-paragraph strips).
            seen_first = False
            for p2 in shape.text_frame.paragraphs:
                if p2 is para:
                    seen_first = True
                    continue
                if seen_first:
                    for r in list(p2.runs):
                        r._r.getparent().remove(r._r)
            return True
    return False


def main(path: str) -> None:
    prs = Presentation(path)
    misses = []
    for idx, old, new in SHAPE_REWRITES:
        if rewrite_shape(prs.slides[idx], old, new):
            print(f"slide {idx + 1}: patched {old[:48]!r}")
        else:
            misses.append((idx, old))
    if misses:
        for idx, old in misses:
            print(f"MISS slide {idx + 1}: {old[:80]!r}")
        raise SystemExit(1)
    prs.save(path)
    print(f"saved {Path(path).name}")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "docs/HireFit_Ranker_Redrob_POLISHED.pptx")
