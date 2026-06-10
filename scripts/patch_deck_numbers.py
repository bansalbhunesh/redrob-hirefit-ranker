"""One-off: patch stale measured numbers inside the polished deck PPTX.

Best-effort run-level text replacement (logs every hit); rebuildable decks
come from build_deck.py, but the POLISHED deck was hand-iterated, so we patch
it in place rather than lose its layout.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

REPLACEMENTS = [
    ("~80–180s", "~3 min worst-case (eval Docker)"),
    ("~80-180s", "~3 min worst-case (eval Docker)"),
    ("80–180s", "177–194s in eval Docker"),
    ("~104s", "≤194s"),
    ("104s", "≤194s"),
    ("123-184 seconds", "177-194 seconds (python:3.11 Docker)"),
    ("123–184 seconds", "177-194 seconds (python:3.11 Docker)"),
    ("123-184s", "177-194s (eval Docker)"),
    ("4.33 GB", "~6.1 GB"),
    ("4.33GB", "~6.1GB"),
    ("auto-pins PYTHONHASHSEED=0 (re-execs once)", "PYTHONHASHSEED=0 pinned by the Dockerfile"),
]


def patch_text_frame(tf, hits):
    for para in tf.paragraphs:
        for run in para.runs:
            t = run.text
            for old, new in REPLACEMENTS:
                if old in t:
                    hits.append((old, t.strip()[:70]))
                    t = t.replace(old, new)
            if t != run.text:
                run.text = t


def walk_shapes(shapes, hits):
    for shape in shapes:
        if shape.shape_type == 6:  # group
            walk_shapes(shape.shapes, hits)
            continue
        if shape.has_text_frame:
            patch_text_frame(shape.text_frame, hits)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    patch_text_frame(cell.text_frame, hits)


def main(path: str) -> None:
    p = Path(path)
    prs = Presentation(str(p))
    hits: list[tuple[str, str]] = []
    for slide in prs.slides:
        walk_shapes(slide.shapes, hits)
    if hits:
        prs.save(str(p))
    print(f"{p.name}: {len(hits)} replacement(s)")
    for old, ctx in hits:
        print(f"  - {old!r} in: {ctx}")


if __name__ == "__main__":
    main(sys.argv[1])
