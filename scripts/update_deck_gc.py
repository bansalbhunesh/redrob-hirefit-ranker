#!/usr/bin/env python3
"""Deck refresh for the GC push (one-off, like patch_deck_numbers.py).

1. Patches stale runtime figures (193s/184s -> audited 187s quiet worst case).
2. Upgrades the VALIDATION slide to the two-judge story.
3. Clones the VALIDATION slide's design into two new slides:
   - PRE-REGISTERED DISCIPLINE (the sweep-we-declined + gold-trap exhibit)
   - INSIDE REDROB (integration + roadmap honesty)
4. Renumbers the hand-written 'n / N' footers.

Usage: python scripts/update_deck_gc.py docs/HireFit_Ranker_Redrob_POLISHED.pptx
"""

from __future__ import annotations

import copy
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

NUM_PATCHES = [
    ("193s", "187s"),
    ("184 s", "187 s"),
    ("in 193s", "in 187s"),
]

# VALIDATION slide text upgrades (exact-match -> replacement).
VALIDATION_EDITS = {
    "An independent LLM judge confirms the top-10 — it isn’t self-graded.":
        "Two independent LLM judges confirm the top-10 — it isn’t self-graded.",
    "Judge   gemini-2.5-flash scores each 0–5 vs the JD rubric: read careers, down-weight unavailable.":
        "Judges   gemini-2.5-flash + gpt-4.1-mini, same 0–5 JD rubric, same 249 ids: kappa 0.935.",
    "Top-10   tiers all 4–5 by the independent judge":
        "Top-10   all 4–5 (judge 1); all tier-5, NDCG@10 = 1.0 (judge 2)",
}

SWEEP_SLIDE = {
    "title": "PRE-REGISTERED DISCIPLINE",
    "headline": "We measured the config that scores higher — and declined it on principle.",
    "sub": "Every knob swept under a decision rule committed before the first run; the JD is the spec, not the judge.",
    "cards": [("0.896", "SHIPPED (LLM-JUDGE)"), ("0.926", "FLOOR-OFF — DECLINED"),
              ("0", "RANK CHANGES 0.25–0.55"), ("11 / 11", "PLANTED TRAPS CAUGHT")],
    "panel1_title": "How the rule decided",
    "panel1_lines": [
        "Rule first   committed before any sweep run — ties break toward the JD's down-weight instruction",
        "Coverage guard   floor 0.85–1.00 configs excluded: <95% label coverage, selection-biased",
        "Winner   shipped config — identical top-100 across floors 0.25–0.55",
    ],
    "panel2_title": "We built our strongest rival — it lost",
    "panel2_body": "A LambdaMART ranker over our own 28 features plus recovered generator structure, trained on 1,500 LLM judgments, evaluated once under a gate committed before training: 0.900 vs our 0.906. Third measured negative result — the hand-tuned ordering stands because nothing we tested beats it, not because nothing was tested.",
    "footer": "Honest framing:  the pool plants 11 perfect-on-paper-but-unavailable seniors (response rates to 0.07) and one honeypot in gold clothing — the guardrails demote all 12; an LLM judge missed one of them.",
}

REDROB_SLIDE = {
    "title": "INSIDE REDROB",
    "headline": "A drop-in scoring layer for Redrob's hiring stack — auditable by design.",
    "sub": "Ranks talent on behavior and evidence, not resumes — the same thesis as the platform it would join.",
    "cards": [("<50 ms", "PER-CANDIDATE AUDIT API"), ("1 cmd", "FULL REPRODUCTION"),
              ("0", "LLM CALLS PER RANKED CANDIDATE"), ("28", "EXPLAINABLE FEATURES")],
    "panel1_title": "Integration path",
    "panel1_lines": [
        "JD compiler   any role text compiles into the same deterministic scoring program",
        "Audit payload   per-candidate JSON: features, multipliers, flags — recruiter-readable",
        "Signals   behavioral multipliers consume Redrob activity data natively",
    ],
    "panel2_title": "Roadmap, honestly",
    "panel2_body": "Indic-script normalization is a localized swap (one module) — the architecture is script-agnostic. When real outcome labels exist, LambdaMART can replace the hand weights over the same 28 features without touching the audit trail: the features are the durable asset.",
    "footer": "LLMs where they are strong — offline judging and evaluation — and deterministic evidence where hiring decisions need to be defended.",
}


def patch_runs(prs):
    hits = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    for old, new in NUM_PATCHES:
                        if old in t:
                            t = t.replace(old, new)
                    if t != run.text:
                        run.text = t
                        hits += 1
    return hits


def set_shape_text(shape, new_text):
    """Replace a shape's text, keeping the first run's formatting."""
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        runs = list(para.runs)
        if not runs:
            continue
        if first:
            runs[0].text = new_text
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
            first = False
        else:
            for r in runs:
                r.text = ""
    return not first


def edit_validation(prs):
    slide = prs.slides[7]
    hits = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        cur = shape.text_frame.text.strip()
        for old, new in VALIDATION_EDITS.items():
            if cur == old and set_shape_text(shape, new):
                hits += 1
    return hits


def clone_slide(prs, src):
    dst = prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(shp._element))
    # Remap image rels so cloned pictures stay valid.
    for rel in list(src.part.rels.values()):
        if "image" in rel.reltype:
            new_rid = dst.part.relate_to(rel.target_part, rel.reltype)
            for blip in dst.shapes._spTree.iter(qn("a:blip")):
                if blip.get(qn("r:embed")) == rel.rId:
                    blip.set(qn("r:embed"), new_rid)
    return dst


def fill_clone(slide, spec):
    """Map the VALIDATION slide's shape roles onto new content by original text."""
    src_map = {
        "VALIDATION": spec["title"],
        "Two independent LLM judges confirm the top-10 — it isn’t self-graded.": spec["headline"],
        "An independent LLM judge confirms the top-10 — it isn’t self-graded.": spec["headline"],
    }
    stats = spec["cards"]
    stat_values = {"0.894": stats[0], "0.873": stats[1], "0.912": stats[2], "1.00": stats[3]}
    stat_labels = {"NDCG@10": stats[0][1], "NDCG@50": stats[1][1],
                   "MAP": stats[2][1], "P@10": stats[3][1]}
    filled = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        cur = shape.text_frame.text.strip()
        if cur in src_map:
            filled += set_shape_text(shape, src_map[cur])
        elif cur.startswith("We avoided grading"):
            filled += set_shape_text(shape, spec["sub"])
        elif cur in stat_values:
            filled += set_shape_text(shape, stat_values[cur][0])
        elif cur in stat_labels:
            filled += set_shape_text(shape, stat_labels[cur])
        elif cur == "How it was judged":
            filled += set_shape_text(shape, spec["panel1_title"])
        elif cur.startswith("249"):
            filled += set_shape_text(shape, spec["panel1_lines"][0])
        elif cur.startswith(("Judge ", "Judges ")):
            filled += set_shape_text(shape, spec["panel1_lines"][1])
        elif cur.startswith("Top-10"):
            filled += set_shape_text(shape, spec["panel1_lines"][2])
        elif cur == "Not propped up by our own labels":
            filled += set_shape_text(shape, spec["panel2_title"])
        elif cur.startswith("The judge composite"):
            filled += set_shape_text(shape, spec["panel2_body"])
        elif cur.startswith("Honest framing"):
            filled += set_shape_text(shape, spec["footer"])
    return filled


def move_slide(prs, old_index, new_index):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[old_index]
    lst.remove(el)
    lst.insert(new_index, el)


def renumber_footers(prs):
    n = len(prs.slides)
    pat = re.compile(r"^\d+ / \d+$")
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame and pat.match(shape.text_frame.text.strip()):
                set_shape_text(shape, f"{i} / {n}")


def main(path: str) -> None:
    p = Path(path)
    prs = Presentation(str(p))
    print(f"number patches: {patch_runs(prs)} runs")
    print(f"validation slide edits: {edit_validation(prs)}")

    src = prs.slides[7]
    s_sweep = clone_slide(prs, src)
    print(f"sweep slide shapes filled: {fill_clone(s_sweep, SWEEP_SLIDE)}")
    s_redrob = clone_slide(prs, src)
    print(f"redrob slide shapes filled: {fill_clone(s_redrob, REDROB_SLIDE)}")

    n = len(prs.slides)
    move_slide(prs, n - 2, 8)    # sweep slide right after VALIDATION
    move_slide(prs, n - 1, 11)   # redrob slide after REPRODUCIBILITY
    renumber_footers(prs)
    prs.save(str(p))
    print(f"saved {p.name} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "docs/HireFit_Ranker_Redrob_POLISHED.pptx")
