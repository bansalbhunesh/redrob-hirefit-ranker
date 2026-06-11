#!/usr/bin/env python3
"""One-off deck fixes after the 2026-06-11 adversarial audit (deck-only; no
ranking artifact is touched).

1. Repairs the PLANTED TRAPS stat card on the PRE-REGISTERED DISCIPLINE slide:
   update_deck_gc.py's renumber_footers() matched the "11 / 11" card value with
   its `^\\d+ / \\d+$` footer pattern and clobbered it to "9 / 13".
2. Unclips the "We built our strongest rival" panel body (12.5pt overflowed the
   1.1in box and hid the last line behind the footer band) by dropping it to 11pt.
3. Copies the per-slide <p:bg> background onto the two cloned slides — the
   actual cause of the washed-out titles in the committed PDF: clone_slide()
   copied shapes but not the slide background, so the near-white headline
   (E6EDF3) rendered on a default-white canvas. Also strips any stray alpha
   (defensive; none found).
4. Retitles THE SYSTEM slide away from "no LLM in the ranking loop" framing.
5. Moves the PRE-REGISTERED DISCIPLINE slide ahead of THE SYSTEM slide and
   fixes the title slide's "see slide 8" footnote (VALIDATION becomes slide 9).
6. Renumbers footers — matching on 9pt font size as well as the `n / N` pattern
   so stat-card values can never be clobbered again.

Usage: python scripts/polish_deck_fixes.py docs/HireFit_Ranker_Redrob_POLISHED.pptx
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

DISCIPLINE_IDX = 8   # PRE-REGISTERED DISCIPLINE (before the move)
REDROB_IDX = 11      # INSIDE REDROB
SYSTEM_IDX = 3       # THE SYSTEM

NEW_SYSTEM_HEADLINE = (
    "Deterministic at rank time — every choice gated by LLM judges offline."
)
NEW_SYSTEM_SUB = (
    "Six CPU-only stages, no network or GPU; tuned and frozen under two LLM "
    "judges and a LambdaMART challenger gate."
)


def set_text_keep_format(shape, new_text: str) -> bool:
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        runs = list(para.runs)
        if not runs:
            continue
        if first:
            runs[0].text = new_text
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
            first = False
        else:
            for r in runs:
                r.text = ""
    return not first


def strip_alpha(shape) -> int:
    removed = 0
    for el in shape._element.iter(qn("a:alpha")):
        el.getparent().remove(el)
        removed += 1
    return removed


def move_slide(prs, old_index: int, new_index: int) -> None:
    lst = prs.slides._sldIdLst
    el = list(lst)[old_index]
    lst.remove(el)
    lst.insert(new_index, el)


def main(path: str) -> None:
    p = Path(path)
    prs = Presentation(str(p))

    # 1. Repair the clobbered stat card (text "9 / 13" at 34pt on the
    #    discipline slide is the damaged card, not a footer).
    fixed_card = 0
    for shape in prs.slides[DISCIPLINE_IDX].shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() != "9 / 13":
            continue
        runs = shape.text_frame.paragraphs[0].runs
        if runs and runs[0].font.size and runs[0].font.size.pt >= 20:
            set_text_keep_format(shape, "11 / 11")
            fixed_card += 1
    print(f"planted-traps card repaired: {fixed_card}")

    # 2. Unclip the rival-panel body.
    unclipped = 0
    for shape in prs.slides[DISCIPLINE_IDX].shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip().startswith("A LambdaMART ranker"):
            for para in shape.text_frame.paragraphs:
                for r in para.runs:
                    r.font.size = Pt(11)
                    unclipped += 1
    print(f"rival-panel runs resized to 11pt: {unclipped}")

    # 3. Copy the clone-source slide's <p:bg> onto both cloned slides (the
    #    washed-title root cause), plus a defensive alpha strip.
    import copy as _copy

    src_bg = prs.slides[7]._element.find(qn("p:cSld")).find(qn("p:bg"))
    bg_fixed = 0
    for idx in (DISCIPLINE_IDX, REDROB_IDX):
        csld = prs.slides[idx]._element.find(qn("p:cSld"))
        if src_bg is not None and csld.find(qn("p:bg")) is None:
            csld.insert(0, _copy.deepcopy(src_bg))
            bg_fixed += 1
    print(f"slide backgrounds copied: {bg_fixed}")
    stripped = 0
    for idx in (DISCIPLINE_IDX, REDROB_IDX):
        for shape in list(prs.slides[idx].shapes)[:3]:
            if shape.has_text_frame:
                stripped += strip_alpha(shape)
    print(f"alpha elements stripped: {stripped}")

    # 4. Retitle THE SYSTEM slide.
    retitled = 0
    for shape in prs.slides[SYSTEM_IDX].shapes:
        if not shape.has_text_frame:
            continue
        cur = shape.text_frame.text.strip()
        if cur.startswith("One offline pipeline scores every candidate"):
            retitled += set_text_keep_format(shape, NEW_SYSTEM_HEADLINE)
        elif cur.startswith("Six deterministic stages turn a raw profile"):
            retitled += set_text_keep_format(shape, NEW_SYSTEM_SUB)
    print(f"system slide shapes retitled: {retitled}")

    # 5. Move the discipline slide ahead of THE SYSTEM slide; VALIDATION
    #    shifts from slide 8 to slide 9, so patch the title footnote.
    move_slide(prs, DISCIPLINE_IDX, SYSTEM_IDX)
    footnote = 0
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for r in para.runs:
                if "see slide 8" in r.text:
                    r.text = r.text.replace("see slide 8", "see slide 9")
                    footnote += 1
    print(f"title footnote patched: {footnote}")

    # 6. Renumber footers — 9pt guard so card values are never matched.
    n = len(prs.slides)
    pat = re.compile(r"^\d+ / \d+$")
    renumbered = 0
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not pat.match(shape.text_frame.text.strip()):
                continue
            runs = shape.text_frame.paragraphs[0].runs
            if runs and runs[0].font.size and runs[0].font.size.pt <= 10:
                set_text_keep_format(shape, f"{i} / {n}")
                renumbered += 1
    print(f"footers renumbered: {renumbered}")

    prs.save(str(p))
    print(f"saved {p.name} with {n} slides")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "docs/HireFit_Ranker_Redrob_POLISHED.pptx")
