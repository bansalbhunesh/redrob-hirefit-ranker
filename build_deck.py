"""Build the submission deck (PPTX) -> convert to PDF separately."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

NAVY = RGBColor(0x0B, 0x12, 0x20)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTE = RGBColor(0x64, 0x74, 0x8B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
CARD = RGBColor(0xF8, 0xFA, 0xFC)

ROOT = Path(__file__).resolve().parent
SS = ROOT / "docs" / "screenshots"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    s = slide.shapes.add_shape(1, 0, 0, W, H)  # rectangle
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element); slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = color
    return tb


def card(slide, l, t, w, h, value, label):
    c = slide.shapes.add_shape(1, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.color.rgb = RGBColor(0xE5,0xE7,0xEB)
    c.line.width = Pt(1); c.shadow.inherit = False
    box(slide, l+Inches(.12), t+Inches(.12), w-Inches(.24), Inches(.6), value, 22, NAVY, bold=True)
    box(slide, l+Inches(.12), t+Inches(.72), w-Inches(.24), Inches(.5), label, 10, MUTE)


def scrim(slide, l, t, w, h, rgb, alpha_pct):
    from pptx.oxml.ns import qn
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb; sh.line.fill.background(); sh.shadow.inherit = False
    srgb = sh.fill.fore_color._xFill.find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))}))
    return sh


def title_slide():
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    cover = SS / "cover.png"
    if cover.exists():
        s.shapes.add_picture(str(cover), 0, 0, width=W, height=H)
        scrim(s, 0, 0, Inches(7.6), H, NAVY, 62)            # left scrim for text contrast
    box(s, Inches(.9), Inches(1.1), Inches(11), Inches(.5),
        "REDROB · INDIA RUNS DATA & AI CHALLENGE", 13, MUTE, bold=True)
    box(s, Inches(.85), Inches(1.7), Inches(11.6), Inches(1.4), "HireFit Ranker", 60, WHITE, bold=True)
    box(s, Inches(.9), Inches(3.1), Inches(11.5), Inches(.7), "Ranks careers, not keywords.", 30, AMBER, bold=True)
    box(s, Inches(.9), Inches(3.95), Inches(11), Inches(1.4),
        "A fast, offline, deterministic engine that surfaces genuinely hireable engineers\n"
        "from a 100,000-candidate pool — a shortlist a recruiter can trust.", 18, LIGHT)
    box(s, Inches(.9), Inches(6.4), Inches(11.5), Inches(.6),
        "github.com/bansalbhunesh/redrob-hirefit-ranker   ·   live demo: redrob-hirefit-ranker.onrender.com",
        13, MUTE)


def header(s, kicker, title, dark=False):
    box(s, Inches(.7), Inches(.5), Inches(12), Inches(.4), kicker, 12,
        AMBER if dark else AMBER, bold=True)
    box(s, Inches(.7), Inches(.95), Inches(12), Inches(.9), title, 30,
        WHITE if dark else NAVY, bold=True)


def problem():
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    header(s, "THE PROBLEM", "The dataset punishes keyword filters")
    box(s, Inches(.7), Inches(2.0), Inches(11.9), Inches(1.0),
        "The JD says \"5–9 yrs, embeddings/retrieval/ranking\" — but it MEANS: find engineers who actually\n"
        "shipped ranking / search / recsys at product companies, even without the buzzwords. And reject:", 16, SLATE)
    items = [
        ("Keyword stuffers", "wrong-role profiles padded with AI skills"),
        ("Honeypots", "impossible profiles — \"expert\" in 10 skills with 0 months used"),
        ("Unavailable", "perfect on paper, but 5% response rate / not open to work"),
    ]
    for i, (t, d) in enumerate(items):
        l = Inches(.7) + i * Inches(4.05)
        c = s.shapes.add_shape(1, l, Inches(3.4), Inches(3.8), Inches(2.2))
        c.fill.solid(); c.fill.fore_color.rgb = CARD; c.line.color.rgb = RGBColor(0xE5,0xE7,0xEB); c.shadow.inherit=False
        box(s, l+Inches(.25), Inches(3.65), Inches(3.4), Inches(.6), t, 19, NAVY, bold=True)
        box(s, l+Inches(.25), Inches(4.35), Inches(3.4), Inches(1.1), d, 14, MUTE)
    box(s, Inches(.7), Inches(6.0), Inches(12), Inches(.8),
        "→ So the ranker reads CAREERS and BEHAVIOR, not skill lists.", 18, AMBER, bold=True)


def approach():
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    header(s, "THE APPROACH", "Hybrid scoring + multiplicative guardrails")
    stages = ["Load","Text","BM25","28-D\nFeatures","Honeypot","Behavioral","Rank","Reason"]
    n = len(stages); gap = Inches(.15); cw = (W - Inches(1.4) - gap*(n-1)) / n
    for i, st in enumerate(stages):
        l = Inches(.7) + i*(cw+gap)
        c = s.shapes.add_shape(1, l, Inches(2.0), cw, Inches(1.0))
        c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background(); c.shadow.inherit=False
        box(s, l, Inches(2.2), cw, Inches(.6), st, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(.7), Inches(3.35), Inches(12), Inches(.7),
        "Final score = base_score  ×  behavioral  ×  honeypot  ×  disqualifier", 18, NAVY, bold=True)
    box(s, Inches(.7), Inches(4.1), Inches(11.9), Inches(2.6),
        "•  BM25 lexical relevance + a 28-feature recruiter matrix (skills, career evidence, seniority, behavior, logistics).\n"
        "•  Guardrails are MULTIPLICATIVE — a high fit score cannot rescue an impossible, off-target, or unavailable profile.\n"
        "•  Career-evidence signals (production, IR/ranking) outweigh skill-list matches → Tier-5s without buzzwords still surface.\n"
        "•  Fully deterministic & offline: same input → byte-identical output; no network, no GPU, no LLM at rank time.", 15, SLATE)


def choices():
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    header(s, "WHY THESE CHOICES", "Deliberate, measured engineering")
    rows = [
        ("Feature matrix + BM25 — not an LLM per candidate",
         "LLM-per-candidate can't scale to 100K on budget, isn't reproducible, needs network. We rank the pool on CPU in ~3 min worst-case in the eval Docker image."),
        ("No dense embeddings — tested & rejected",
         "Built a model2vec/potion branch, gated on a measured A/B: NDCG@10 +0.0000, ~2.2× runtime → FAIL. Shipped the simpler system."),
        ("Career-evidence over keywords",
         "Production / IR-ranking signals from career history outweigh skill lists — how non-buzzword Tier-5s surface."),
        ("Multiplicative behavioral & honeypot guardrails",
         "Encodes the JD's explicit \"down-weight the unavailable\" rule; impossible profiles forced to zero."),
        ("Deterministic everything",
         "Byte-identical, traceable, debuggable — the basis for explainability and reproduction."),
    ]
    y = Inches(1.95)
    for t, d in rows:
        box(s, Inches(.7), y, Inches(5.0), Inches(.9), t, 14, NAVY, bold=True)
        box(s, Inches(5.9), y, Inches(6.7), Inches(.9), d, 13, SLATE)
        y += Inches(1.0)


def results():
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    header(s, "RESULTS", "A shortlist a recruiter can trust")
    kpis = [("≤194s","100K in eval Docker · CPU"),("100% offline","no network · GPU · LLM"),
            ("0 / 53","honeypots in top output"),("P@10 · 1.0","independent LLM-judge")]
    cw = Inches(2.9)
    for i,(v,l) in enumerate(kpis):
        card(s, Inches(.7)+i*Inches(3.05), Inches(1.9), cw, Inches(1.35), v, l)
    img = SS / "dashboard.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(.7), Inches(3.5), width=Inches(11.9))
    box(s, Inches(.7), Inches(7.0), Inches(12), Inches(.4),
        "Live interactive dashboard — pipeline stages, honeypot blocking, per-candidate audit.", 12, MUTE)


def validation():
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    header(s, "VALIDATION", "Proven without the hidden labels", dark=True)
    box(s, Inches(.7), Inches(2.0), Inches(11.9), Inches(1.2),
        "1.  Non-circular heuristic eval — an independent labeler sharing NO code with the ranker → composite 0.881\n"
        "     (rules out self-grading).", 16, LIGHT)
    box(s, Inches(.7), Inches(3.3), Inches(11.9), Inches(1.6),
        "2.  LLM-as-judge (dev-only, never in the ranking path) scored the top-10 against the JD:", 16, LIGHT)
    box(s, Inches(1.1), Inches(4.1), Inches(11), Inches(.8),
        "top-10 tiers [5,5,4,4,5,5,5,5,5,5]   ·   P@10 = 1.0   ·   NDCG@10 = 0.894", 20, AMBER, bold=True)
    box(s, Inches(.7), Inches(5.2), Inches(11.9), Inches(1.6),
        "The judge even confirmed our behavioral guardrails are MORE recruiter-aware than itself —\n"
        "we correctly down-weighted high-skill candidates with 12% response rate that the LLM over-rated.\n\n"
        "\"We don't rank resumes. We rank hireable candidates.\"", 16, LIGHT)


def explain():
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    header(s, "EXPLAINABILITY", "Every score is traceable — nothing is a black box")
    box(s, Inches(.7), Inches(1.9), Inches(6.0), Inches(4.5),
        "•  Per-candidate REASONING in the output, generated only from facts in the profile — no hallucinated skills.\n\n"
        "•  Every score DECOMPOSES into named features + the three multipliers.\n\n"
        "•  Interactive dashboard shows the pipeline, honeypot blocking, and a live feature + reasoning audit.\n\n"
        "•  Two live UIs: Render dashboard + HuggingFace sandbox.", 16, SLATE)
    img = SS / "space-results.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(7.0), Inches(2.0), width=Inches(5.6))


def enterprise():
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    header(s, "ENTERPRISE UPGRADE", "Lab Architecture vs Baseline Prototype")
    box(s, Inches(.7), Inches(1.9), Inches(11.9), Inches(4.5),
        "•  TEST COVERAGE (115 ➔ 163): Expanded coverage by +42%, including 13 rigorous counterfactual tests bounding location, gender, and name proxies mathematically.\n\n"
        "•  BACKEND STORAGE (Dicts ➔ SQLite): Migrated from volatile in-memory storage to a persistent SQLite WAL database, ensuring batch operations survive server restarts.\n\n"
        "•  API SECURITY (Open ➔ X-Demo-Token): Guarded all write-endpoints with token authentication to prevent unauthorized abuse in production.\n\n"
        "•  SYSTEM RATING (89/100 ➔ 91/100 self-assessed; independent forensic audit ~88/100): Closed the 100-score gap with the official ranking now byte-reproducible in the pinned Docker image; ~80s on a cloud 2-vCPU runner (~130-155s in local Docker); zero honeypots in the top 100.", 16, SLATE)


def close():
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    box(s, Inches(.9), Inches(1.3), Inches(11.5), Inches(1.0), "Submission", 40, WHITE, bold=True)
    box(s, Inches(.9), Inches(2.6), Inches(11.5), Inches(2.6),
        "✓  GitHub repo — clean, complete, working code\n"
        "✓  Ranked output — submission.csv (candidate_id, rank, score, reasoning)\n"
        "✓  Methodology — METHODOLOGY.md + this deck\n"
        "✓  Live demos — Render dashboard + HuggingFace Space", 20, LIGHT)
    box(s, Inches(.9), Inches(5.6), Inches(11.5), Inches(1.0),
        "Reproduce:  python rank.py --candidates ./candidates.jsonl --out ./submission.csv", 16, AMBER, bold=True)
    box(s, Inches(.9), Inches(6.6), Inches(11.5), Inches(.5),
        "github.com/bansalbhunesh/redrob-hirefit-ranker", 13, MUTE)


title_slide(); problem(); approach(); choices(); results(); validation(); explain(); enterprise(); close()
out = ROOT / "docs" / "Redrob_HireFit_Ranker.pptx"
prs.save(str(out))
print("saved", out, "| slides:", len(prs.slides._sldIdLst))
